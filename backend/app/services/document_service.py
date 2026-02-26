"""
document_service.py — Extract text from uploaded documents and split into sentences.
Supports PDF, DOCX, and TXT.
"""
import os
import re
from typing import List


def extract_text(file_path: str) -> str:
    """Extract all text from a document file. Returns empty string on failure."""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        elif ext in (".docx", ".doc"):
            return _extract_docx(file_path)
        elif ext in (".pptx", ".ppt"):
            return _extract_pptx(file_path)
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        else:
            # Try reading as plain text for unknown types
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
    except Exception as e:
        print(f"[DocService] Failed to extract text from {file_path}: {e}")
        return ""


def _extract_pdf(file_path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n".join(pages)


def _extract_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_runs.append(shape.text)
    return "\n".join(text_runs)


def split_sentences(text: str, min_len: int = 15) -> List[str]:
    """Split text into sentences; discard very short fragments."""
    if not text:
        return []
    # Split on sentence-ending punctuation or newlines
    raw = re.split(r'(?<=[.!?])\s+|\n{2,}', text)
    sentences = []
    for s in raw:
        s = s.strip()
        if len(s) >= min_len:
            sentences.append(s)
    return sentences
